
import tempfile
from pathlib import Path
import re
from datetime import date
import streamlit as st
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def generate_new_pptx_option1(analysis_data, title_bg_info, second_page_bg_info, content_bg_info):
    """Generate a PPTX presentation from analysis_data."""
    prs = Presentation()
    blank = prs.slide_layouts[-1]

    def apply_bg(slide, image_bytes, image_name, default_color):
        w, h = prs.slide_width, prs.slide_height
        if image_bytes:
            try:
                with tempfile.TemporaryDirectory() as td:
                    pth = Path(td) / image_name
                    with open(pth, "wb") as f: f.write(image_bytes)
                    slide.background.fill.background()
                    slide.shapes.add_picture(str(pth), 0, 0, w, h)
            except Exception as e:
                st.warning(f"Could not set background image ({image_name}): {e}. Using solid fill.")
                slide.background.fill.solid(); slide.background.fill.fore_color.rgb = default_color
        else:
            slide.background.fill.solid(); slide.background.fill.fore_color.rgb = default_color

    # Title slide
    s0 = prs.slides.add_slide(blank)
    apply_bg(s0, title_bg_info.get("bytes"), title_bg_info.get("name"), RGBColor(255,255,255))
    bg = analysis_data.get("Prospect & RFP Background", {})
    client = "N/A"; rfp_no = "N/A"; proj = "Project Name"
    if isinstance(bg, dict): client = bg.get("Client Name", client)
    if isinstance(bg, str):
        m = re.search(r"(?:Client Name|Client):\s*(.*?)(?:\n|$)", bg, re.I)
        if m: client = m.group(1).strip()
        m2 = re.search(r"(?:RFP Number|RFQ Number):\s*(.*?)(?:\n|$)", bg, re.I)
        if m2: rfp_no = m2.group(1).strip()
        m3 = re.search(r"(?:Project Name|Project Title):\s*(.*?)(?:\n|$)", bg, re.I)
        if m3 and m3.group(1).strip(): proj = m3.group(1).strip()
    today = date.today().strftime("%B %d, %Y")

    tb = s0.shapes.add_textbox(Inches(1), Inches(2.5), prs.slide_width - Inches(2), Inches(5)).text_frame
    tb.clear()
    def p(text, size=18, bold=False, color=RGBColor(255,255,255)):
        para = tb.add_paragraph(); para.text = text
        para.font.size = Pt(size); para.font.bold = bold; para.font.color.rgb = color
        para.alignment = PP_ALIGN.LEFT
    p(client, 22, True); p(f"RFP Number: {rfp_no}", 18); p(""); p(proj, 18); p(""); p(today, 18)

    # Sections
    sections = [
        "Prospect & RFP Background",
        "Scope of Work Details",
        "Service Level Agreements (SLAs)",
        "RFP Submission Information",
        "RFP Schedule",
        "RFP Evaluation Criteria",
        "SWOT Analysis",
    ]

    def section_slide(title, content, center=False):
        s = prs.slides.add_slide(blank)
        apply_bg(s, content_bg_info.get("bytes"), content_bg_info.get("name"), RGBColor(0,0,0))
        head = s.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width-Inches(1), Inches(0.7)).text_frame
        head.text = title; head.paragraphs[0].font.size = Pt(24); head.paragraphs[0].font.bold = True
        head.paragraphs[0].font.color.rgb = RGBColor(0x03, 0x9e, 0xed)

        body = s.shapes.add_textbox(Inches(0.5), Inches(1.2), prs.slide_width-Inches(1), prs.slide_height-Inches(1.7)).text_frame
        body.clear(); body.word_wrap = True
        if isinstance(content, dict):
            for k, v in content.items():
                pr = body.add_paragraph(); pr.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
                runk = pr.add_run(); runk.text = f"{k.upper()}: "; runk.font.bold = True; runk.font.size = Pt(16)
                if isinstance(v, list):
                    for item in v:
                        bp = body.add_paragraph(); bp.text = str(item); bp.level = 1; bp.font.size = Pt(16)
                else:
                    rv = pr.add_run(); rv.text = str(v); rv.font.size = Pt(16)
        else:
            lines = (content or "").split("\\n")
            for line in lines:
                if line.strip():
                    pr = body.add_paragraph(); pr.text = line.strip(); pr.font.size = Pt(16)
                    pr.alignment = PP_ALIGN.CENTER if center else PP_ALIGN.LEFT
        return s

    for sec in sections:
        c = analysis_data.get(sec)
        if not c or (isinstance(c, (dict, list)) and not c): continue
        center = (sec == "RFP Schedule")
        section_slide(sec, c, center=center)

    buf = BytesIO(); prs.save(buf); buf.seek(0)
    return buf
