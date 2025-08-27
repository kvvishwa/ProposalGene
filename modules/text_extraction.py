# modules/text_extraction.py
# -----------------------------------------------------------------------------
# Text extraction helpers with optional OCR for PDFs.
# - extract_text(path): returns text best-effort (no OCR by default)
# - ocr_pdf(path, max_pages=20, dpi=250): OCR fallback (guarded; returns "" if deps missing)
# -----------------------------------------------------------------------------

from __future__ import annotations
import os
from typing import Optional

# You likely already use these; keep your existing imports for PDF/DOCX/PPTX
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None

try:
    import docx2txt
except Exception:
    docx2txt = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

# Optional OCR deps (guarded)
try:
    from pdf2image import convert_from_path
    import pytesseract
except Exception:
    convert_from_path = None
    pytesseract = None


def _pdf_text(path: str) -> str:
    if fitz is None:
        return ""
    try:
        doc = fitz.open(path)
        texts = []
        for page in doc:
            texts.append(page.get_text("text") or "")
        doc.close()
        return "\n".join(texts).strip()
    except Exception:
        return ""


def _docx_text(path: str) -> str:
    if docx2txt is None:
        return ""
    try:
        return (docx2txt.process(path) or "").strip()
    except Exception:
        return ""


def _pptx_text(path: str) -> str:
    if Presentation is None:
        return ""
    try:
        prs = Presentation(path)
        chunks = []
        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, "text") and shp.text:
                    chunks.append(shp.text)
        return "\n".join(chunks).strip()
    except Exception:
        return ""


def extract_text(path: str) -> str:
    """
    Best-effort text extraction (no OCR). Use ocr_pdf(path) explicitly if needed.
    """
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return _pdf_text(path)
    if ext == ".docx":
        return _docx_text(path)
    if ext == ".pptx":
        return _pptx_text(path)
    # simple pass-through for text files
    if ext in (".txt",):
        try:
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()
        except Exception:
            return ""
    return ""


def ocr_pdf(path: str, max_pages: int = 20, dpi: int = 250) -> str:
    """
    OCR fallback for scanned PDFs. Requires `pdf2image` and `pytesseract`.
    Returns "" if dependencies are not installed.
    """
    if convert_from_path is None or pytesseract is None:
        return ""
    text_chunks = []
    try:
        pages = convert_from_path(path, dpi=dpi)
        for i, img in enumerate(pages):
            if i >= max_pages:
                break
            try:
                text_chunks.append(pytesseract.image_to_string(img))
            except Exception:
                continue
        return "\n".join(text_chunks).strip()
    except Exception:
        return ""
